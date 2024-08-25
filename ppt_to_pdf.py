from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from PIL import Image
import os


def ppt_to_pdf(input_file, output_file=None):
    """
    Converts a PowerPoint file to PDF by rendering each slide as an image and embedding in a PDF.

    :param input_file: Path to the input PowerPoint file (.pptx).
    :param output_file: Path to save the output PDF file. If None, saves in the same location as the input file with a .pdf extension.
    """
    # Check if the input file exists
    if not os.path.exists(input_file):
        raise FileNotFoundError(f"The file {input_file} does not exist.")

    # Set the output file path
    if output_file is None:
        output_file = os.path.splitext(input_file)[0] + ".pdf"

    # Load the presentation
    prs = Presentation(input_file)

    # Create a PDF canvas
    pdf_canvas = canvas.Canvas(output_file, pagesize=letter)

    # Iterate over slides
    for idx, slide in enumerate(prs.slides):
        slide_image = f"slide_{idx}.png"
        # Save the slide as an image
        slide.shapes._spTree.export(slide_image)

        # Open image and add it to the PDF
        img = Image.open(slide_image)
        img_width, img_height = img.size

        # Scale image to fit the page
        scale_factor = min(letter[0] / img_width, letter[1] / img_height)
        pdf_canvas.drawImage(
            slide_image, 0, 0, img_width * scale_factor, img_height * scale_factor
        )

        # Move to the next page in PDF
        pdf_canvas.showPage()

        # Remove the image file
        os.remove(slide_image)

    # Save the PDF
    pdf_canvas.save()

    return output_file
